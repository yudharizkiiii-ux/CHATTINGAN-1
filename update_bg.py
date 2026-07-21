from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\user\Documents\CHATTING\Portfolio_Antigravity_v3.pptx"
OUT = r"C:\Users\user\Documents\CHATTING\Portfolio_Antigravity_v4.pptx"

# Gradient palettes per slide (1-indexed): (stop1, stop2, stop3, angle_deg)
# Base = deep navy; accent shifts per section for variety but cohesive
PALETTES = {
    1:  ("0A1428", "0E2A4A", "00A884", 45),   # title: navy -> teal
    2:  ("0A1428", "1A1B4B", "6D2E8E", 60),   # agenda: navy -> purple
    3:  ("0A1428", "0E3B5C", "00A884", 50),   # gamehub overview: navy -> teal
    4:  ("0A1428", "0E3B5C", "005C4B", 50),   # chat overview: navy -> emerald
    5:  ("0A1428", "221B5C", "3B2E8E", 55),   # tech stack: navy -> indigo
    6:  ("0A1428", "3B2A0E", "B8860B", 45),   # gamehub steps: navy -> amber
    7:  ("0A1428", "3B0E2A", "8E2E5C", 55),   # chat steps: navy -> rose
    8:  ("0A1428", "0E3B5C", "02C39A", 50),   # gamehub features: navy -> mint
    9:  ("0A1428", "0E3B5C", "005C4B", 50),   # chat features: navy -> emerald
    10: ("0A1428", "1A1B4B", "6D2E8E", 60),   # how to run: navy -> purple
    11: ("0A1428", "3B2A0E", "F9E795", 45),   # 5-day experience: navy -> gold
    12: ("0A1428", "1A1B2B", "065A82", 50),   # screenshots: navy -> ocean
    13: ("0A1428", "0E2A4A", "00A884", 45),   # closing: navy -> teal
}

p = Presentation(SRC)

def set_gradient(shape, c1, c2, c3, angle_deg):
    """Replace shape fill with a 3-stop linear gradient."""
    sp = shape._element
    # remove existing <a:gradFill>, <a:solidFill>, <a:noFill>, etc under spPr
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        # try alternate namespace for autoshapes
        spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return False
    # remove all fill children
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    # build new gradient XML
    # angle is in 60000ths of a degree
    ang = str(int(angle_deg * 60000))
    grad_xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
        f'<a:gs pos="50000"><a:srgbClr val="{c2}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c3}"/></a:gs>'
        f'</a:gsLst>'
        f'<a:lin ang="{ang}" scaled="1"/>'
        '</a:gradFill>'
    )
    grad_el = etree.fromstring(grad_xml)
    # insert after <a:prstGeom> (or wherever fill goes in spPr sequence)
    # Order in CT_ShapeProperties: xfrm, prstGeom/custGeom, noFill/solidFill/gradFill/..., ln, effectLst, scene3d, sp3d, extLst
    prstGeom = spPr.find(qn("a:prstGeom"))
    ln = spPr.find(qn("a:ln"))
    if prstGeom is not None:
        prstGeom.addnext(grad_el)
    elif ln is not None:
        ln.addprevious(grad_el)
    else:
        spPr.insert(0, grad_el)
    return True

count = 0
for i, slide in enumerate(p.slides, 1):
    shapes = list(slide.shapes)
    if not shapes:
        continue
    bg_shape = shapes[0]
    # only update the full-slide bg rectangle
    if bg_shape.left == 0 and bg_shape.top == 0 and abs(bg_shape.width - p.slide_width) < 100000:
        c1, c2, c3, ang = PALETTES.get(i, PALETTES[1])
        if set_gradient(bg_shape, c1, c2, c3, ang):
            count += 1
            print(f"  slide {i}: gradient applied ({c1}->{c2}->{c3} @{ang}deg)")
        else:
            print(f"  slide {i}: FAILED to apply gradient")
    else:
        print(f"  slide {i}: shape0 not bg, skip")

p.save(OUT)
print(f"\nSaved {OUT}")
print(f"Backgrounds updated on {count}/13 slides")

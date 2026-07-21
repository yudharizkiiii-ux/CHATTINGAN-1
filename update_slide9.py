from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

SRC = r"C:\Users\user\Documents\CHATTING\Portfolio_Antigravity_v4.pptx"
OUT = r"C:\Users\user\Documents\CHATTING\Portfolio_Antigravity_v5.pptx"
PIC_DIR = r"C:\Users\user\Pictures"

# Foto fitur + caption (urutan untuk 5 kolom horizontal)
FOTOS = [
    ("FITUR CHAT.PNG",    "Pesan Instan"),
    ("FITUR TELEPON.PNG", "Voice Call"),
    ("FITUR VC.PNG",      "Video Call"),
    ("FITUR VN.PNG",      "Voice Notes"),
    ("KIRIM GAMBAR.PNG",  "Kirim Gambar"),
]

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0xA8, 0x84)
LIGHT = RGBColor(0xE9, 0xED, 0xE9)
GRAY = RGBColor(0x86, 0x96, 0x87)

p = Presentation(SRC)
slide = p.slides[8]  # slide 9 (0-indexed)

# 1. Hapus semua shape kecuali 3 pertama (bg, judul, subtitle)
#    Shape index 3..32 = 6 card fitur (5 shapes per card)
shapes_to_remove = []
for i, sh in enumerate(slide.shapes):
    if i >= 3:  # keep shape 0 (bg), 1 (title), 2 (subtitle)
        shapes_to_remove.append(sh)
for sh in shapes_to_remove:
    sp = sh._element
    sp.getparent().remove(sp)
print(f"Removed {len(shapes_to_remove)} shapes from slide 9")

# 2. Tambah 5 foto fitur dalam baris horizontal di tengah
# Layout: 5 foto, masing-masing 2.3" wide, 3.6" tall, gap 0.15"
# Foto height 3.2" + caption 0.4"
slide_w = p.slide_width  # EMU
slide_h = p.slide_height

# Konversi ke inches untuk perhitungan (1 inch = 914400 EMU)
fw = Inches(2.3)       # foto width
fh = Inches(3.2)      # foto height
gap = Inches(0.15)    # gap antar foto
total_w = 5 * 914400 * 2.3 + 4 * 914400 * 0.15  # EMU
start_x = (slide_w - total_w) // 2
y_foto = Inches(1.95)   # y untuk foto
y_caption = Inches(5.25) # y untuk caption (di bawah foto)

print(f"Total width: {total_w/914400:.2f}\"  start_x: {start_x/914400:.2f}\"")

for idx, (fname, caption) in enumerate(FOTOS):
    fpath = os.path.join(PIC_DIR, fname)
    x = start_x + idx * (fw + gap)
    if not os.path.exists(fpath):
        print(f"  SKIP {fname} (not found)")
        continue
    # Tambah foto dengan sizing contain agar rasio tidak terdistorsi
    pic = slide.shapes.add_picture(fpath, x, y_foto, fw, fh)
    # Tambah caption di bawah foto
    tb = slide.shapes.add_textbox(x, y_caption, fw, Inches(0.4))
    tf = tb.text_frame
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = caption
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = "Arial"
    print(f"  Added {fname} -> caption '{caption}' at x={x/914400:.2f}\"")

p.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Slide 9 now has {len(p.slides[8].shapes)} shapes")

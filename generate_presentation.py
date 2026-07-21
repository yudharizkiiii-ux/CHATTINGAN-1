"""
Generate PPTX Presentasi: MQTT WhatsApp Web Clone
Aplikasi Real-Time Chat dengan MQTT + WebRTC
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.util as util

# ─── PALET WARNA ──────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0b, 0x14, 0x1a)   # #0b141a  – latar gelap WhatsApp
BG_CARD       = RGBColor(0x1f, 0x2c, 0x34)   # #1f2c34  – card/sidebar
ACCENT_GREEN  = RGBColor(0x00, 0xa8, 0x84)   # #00a884  – hijau WhatsApp
ACCENT_TEAL   = RGBColor(0x00, 0x5c, 0x4b)   # #005c4b  – teal WhatsApp
ACCENT_BLUE   = RGBColor(0x53, 0xbd, 0xeb)   # #53bdeb  – biru info
WHITE         = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY    = RGBColor(0xc0, 0xc9, 0xd0)   # teks sekunder
YELLOW        = RGBColor(0xff, 0xd7, 0x4e)   # aksen kuning
ORANGE        = RGBColor(0xff, 0x6b, 0x35)   # aksen oranye

MSO_RECT = 1
MSO_OVAL = 9


def set_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0), shape_type=MSO_RECT):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=Pt(18), bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_accent_bar(slide, left=0, top=0, width=10, height=0.06, color=ACCENT_GREEN):
    bar = slide.shapes.add_shape(MSO_RECT,
        Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


# ─── SLIDE 1: COVER ───────────────────────────────────────────────────────────
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    card = add_rect(slide, 0.5, 1.2, 9, 5.5, fill_color=BG_CARD, line_color=ACCENT_GREEN, line_width=Pt(1.5))
    add_accent_bar(slide, left=0.5, top=1.2, width=9, height=0.07, color=ACCENT_GREEN)

    icon_circle = slide.shapes.add_shape(MSO_OVAL,
        Inches(3.8), Inches(1.6), Inches(2.4), Inches(2.4))
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = ACCENT_GREEN
    icon_circle.line.fill.background()

    add_text_box(slide, "💬", left=4.1, top=2.0, width=1.8, height=1.5,
                 font_size=Pt(54), align=PP_ALIGN.CENTER)

    add_text_box(slide, "MQTT WhatsApp Clone",
                 left=0.5, top=3.95, width=9, height=0.7,
                 font_size=Pt(34), bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, "Aplikasi Chat Real-Time | MQTT + WebRTC",
                 left=0.5, top=4.6, width=9, height=0.5,
                 font_size=Pt(16), bold=False, color=ACCENT_GREEN,
                 align=PP_ALIGN.CENTER)

    add_accent_bar(slide, left=2.5, top=5.1, width=5, height=0.04, color=ACCENT_TEAL)

    add_text_box(slide, "React + Vite  •  MQTT.js  •  WebRTC  •  Tailwind CSS",
                 left=0.5, top=5.3, width=9, height=0.4,
                 font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "2026",
                 left=0.5, top=5.7, width=9, height=0.3,
                 font_size=Pt(10), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_accent_bar(slide, left=0, top=6.8, width=10, height=0.08, color=ACCENT_GREEN)


# ─── SLIDE 2: OVERVIEW PROYEK ─────────────────────────────────────────────────
def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "📋  Overview Proyek", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "Apa itu MQTT WhatsApp Clone?", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    desc_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(9.2), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("Aplikasi web chat real-time premium berbasis dark-theme glassmorphism, "
                "terinspirasi dari WhatsApp Web. Dibangun menggunakan protokol MQTT lewat "
                "WebSocket untuk pengiriman pesan instan, dan WebRTC untuk panggilan suara "
                "& video langsung antar browser (Peer-to-Peer) tanpa server relay.")
    run.font.size = Pt(13)
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = "Segoe UI"

    stats = [
        ("⚡", "< 1 detik", "Latensi Pesan"),
        ("🖼️", "10 MB", "Maks Ukuran Foto"),
        ("🎙️", "250 KB", "Maks Voice Note"),
        ("🌐", "P2P", "Video Call"),
    ]
    card_w, card_h = 2.1, 1.8
    for i, (icon, value, label) in enumerate(stats):
        x = 0.35 + i * 2.35
        add_rect(slide, x, 2.95, card_w, card_h, fill_color=BG_CARD, line_color=ACCENT_GREEN, line_width=Pt(1))
        add_text_box(slide, icon, left=x, top=3.0, width=card_w, height=0.5,
                     font_size=Pt(22), align=PP_ALIGN.CENTER)
        add_text_box(slide, value, left=x, top=3.5, width=card_w, height=0.55,
                     font_size=Pt(22), bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, left=x, top=4.05, width=card_w, height=0.5,
                     font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 5.1, 10, 0.8, fill_color=ACCENT_TEAL)
    add_text_box(slide, "🛠️  Stack:   React 18  •  Vite  •  MQTT.js 5  •  WebRTC  •  Tailwind CSS  •  Canvas API  •  MediaRecorder API",
                 left=0.3, top=5.18, width=9.5, height=0.55,
                 font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 3: FITUR UTAMA ─────────────────────────────────────────────────────
def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "✨  Fitur-Fitur Unggulan", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "6 Fitur Utama yang Membuat Aplikasi Ini Premium", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    features = [
        ("📡", "Pesan Real-Time (MQTT)", "Kirim & terima pesan dalam hitungan milidetik menggunakan protokol MQTT lewat WSS (WebSocket Secure).", ACCENT_GREEN),
        ("📹", "Voice & Video Call (WebRTC)", "Panggilan suara dan video langsung P2P antar browser. MQTT dipakai sebagai signaling channel untuk SDP & ICE.", ACCENT_BLUE),
        ("🎙️", "Voice Notes", "Rekam audio langsung dari browser via MediaRecorder API. Dikirim terkompresi Base64 WebM (<250 KB).", YELLOW),
        ("🖼️", "Kompresi Foto 10MB", "Kirim foto besar tanpa hambatan. Kompres sisi klien via Canvas API memastikan ukuran <150 KB.", ORANGE),
        ("👁️", "Presence & Typing Indicator", "Pantau siapa saja yang online di grup secara real-time, plus indikator 'sedang mengetik...' yang muncul dinamis.", ACCENT_GREEN),
        ("🌙", "Desain Glassmorphism", "Dark mode premium bergaya WhatsApp Web dengan efek blur kaca, gradien mulus, dan animasi transisi elegan.", ACCENT_BLUE),
    ]

    col_w, card_h = 4.5, 1.35
    for i, (icon, title, desc, color) in enumerate(features):
        col = i % 2
        row = i // 2
        x = 0.25 + col * 4.85
        y = 1.35 + row * 1.5

        add_rect(slide, x, y, col_w, card_h, fill_color=BG_CARD, line_color=color, line_width=Pt(1.2))
        add_rect(slide, x, y, 0.07, card_h, fill_color=color)

        add_text_box(slide, icon + "  " + title,
                     left=x+0.15, top=y+0.05, width=col_w-0.2, height=0.4,
                     font_size=Pt(13), bold=True, color=color)
        add_text_box(slide, desc,
                     left=x+0.15, top=y+0.42, width=col_w-0.25, height=0.85,
                     font_size=Pt(10.5), color=LIGHT_GRAY)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 4: ARSITEKTUR SISTEM ───────────────────────────────────────────────
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "🏗️  Arsitektur Sistem", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "Cara kerja komunikasi antar komponen", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    add_rect(slide, 0.3, 1.4, 2.5, 2.8, fill_color=RGBColor(0x12, 0x3a, 0x2f), line_color=ACCENT_GREEN, line_width=Pt(1.5))
    add_text_box(slide, "🖥️  Browser A\n(Pengirim)", left=0.3, top=1.4, width=2.5, height=0.7,
                 font_size=Pt(12), bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, "React App\n+ MQTT Client\n+ WebRTC Peer",
                 left=0.3, top=2.0, width=2.5, height=1.2,
                 font_size=Pt(10), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 3.6, 1.4, 2.8, 2.8, fill_color=RGBColor(0x1a, 0x1a, 0x4e), line_color=ACCENT_BLUE, line_width=Pt(2))
    add_text_box(slide, "☁️  HiveMQ Broker\n(Cloud)", left=3.6, top=1.4, width=2.8, height=0.75,
                 font_size=Pt(13), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "wss://broker.hivemq.com\n:8884/mqtt\n\nTopik MQTT:\n• mqtt_chat/messages/{room}\n• mqtt_chat/presence\n• mqtt_chat/typing\n• mqtt_chat/calls/{room}",
                 left=3.6, top=2.1, width=2.8, height=2.0,
                 font_size=Pt(9), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 7.2, 1.4, 2.5, 2.8, fill_color=RGBColor(0x12, 0x3a, 0x2f), line_color=ACCENT_GREEN, line_width=Pt(1.5))
    add_text_box(slide, "🖥️  Browser B\n(Penerima)", left=7.2, top=1.4, width=2.5, height=0.7,
                 font_size=Pt(12), bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, "React App\n+ MQTT Client\n+ WebRTC Peer",
                 left=7.2, top=2.0, width=2.5, height=1.2,
                 font_size=Pt(10), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "--- MQTT MSG -->", left=2.75, top=2.6, width=0.9, height=0.4,
                 font_size=Pt(9), color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, "<-- MQTT MSG ---", left=6.35, top=2.6, width=0.9, height=0.4,
                 font_size=Pt(9), color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    add_rect(slide, 2.0, 4.55, 6.0, 0.55, fill_color=RGBColor(0x2d, 0x1a, 0x00), line_color=YELLOW, line_width=Pt(1))
    add_text_box(slide, "🎥  WebRTC P2P  --  Audio / Video Stream Langsung  --  🎥",
                 left=2.0, top=4.58, width=6.0, height=0.45,
                 font_size=Pt(9.5), color=YELLOW, align=PP_ALIGN.CENTER, bold=True)

    add_text_box(slide, "💡 Signaling WebRTC (SDP & ICE) lewat MQTT, lalu stream langsung P2P via STUN Google",
                 left=0.3, top=5.25, width=9.5, height=0.5,
                 font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 5: PROSES DEVELOPMENT ─────────────────────────────────────────────
def slide_development(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "🛠️  Proses Pengembangan", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "Step-by-Step dari Nol Sampai Jadi", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    steps = [
        ("1", "Inisialisasi Proyek", "npm create vite@latest CHATTING -- --template react && npm install", ACCENT_GREEN),
        ("2", "Setup Tailwind CSS", "npm install -D tailwindcss postcss autoprefixer  |  Konfigurasi palet warna WhatsApp gelap", ACCENT_BLUE),
        ("3", "Install Dependensi Utama", "npm install mqtt lucide-react canvas-confetti  +  aedes websocket-stream", YELLOW),
        ("4", "Bangun App.jsx (MQTT & WebRTC)", "Koneksi HiveMQ WSS, presence/typing listener, alur WebRTC P2P dengan SDP & ICE handshake", ORANGE),
        ("5", "Versi Standalone HTML", "Port seluruh logika ke mqtt-chat-easy.html  |  Semua library via CDN, zero install", ACCENT_GREEN),
        ("6", "Git Init & Commit", "git init  |  git add .  |  git commit -m 'Initial commit: WhatsApp Chat over MQTT'", ACCENT_BLUE),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        y = 1.35 + i * 0.9

        circle = slide.shapes.add_shape(MSO_OVAL, Inches(0.3), Inches(y+0.05), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text_box(slide, num, left=0.3, top=y+0.02, width=0.55, height=0.55,
                     font_size=Pt(13), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            line_bar = slide.shapes.add_shape(MSO_RECT,
                Inches(0.54), Inches(y+0.6), Inches(0.07), Inches(0.33))
            line_bar.fill.solid()
            line_bar.fill.fore_color.rgb = color
            line_bar.line.fill.background()

        add_rect(slide, 1.05, y, 8.65, 0.75, fill_color=BG_CARD, line_color=color, line_width=Pt(0.8))
        add_text_box(slide, title, left=1.2, top=y+0.02, width=3.0, height=0.35,
                     font_size=Pt(12), bold=True, color=color)
        add_text_box(slide, desc, left=1.2, top=y+0.35, width=8.3, height=0.4,
                     font_size=Pt(9.5), color=LIGHT_GRAY)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 6: CARA MENJALANKAN ────────────────────────────────────────────────
def slide_howto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "🚀  Cara Menjalankan Aplikasi", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "Dua opsi fleksibel — pilih sesuai kebutuhan", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    # Opsi A
    add_rect(slide, 0.3, 1.35, 4.4, 3.8, fill_color=BG_CARD, line_color=ACCENT_GREEN, line_width=Pt(1.5))
    add_accent_bar(slide, 0.3, 1.35, 4.4, 0.06, ACCENT_GREEN)
    add_text_box(slide, "⚙️  Opsi A — React + Vite", left=0.4, top=1.42, width=4.2, height=0.55,
                 font_size=Pt(14), bold=True, color=ACCENT_GREEN)
    add_text_box(slide, "Rekomendasi untuk pengembangan & debug",
                 left=0.4, top=1.9, width=4.1, height=0.35,
                 font_size=Pt(10), color=LIGHT_GRAY)

    add_rect(slide, 0.45, 2.3, 4.1, 1.5, fill_color=RGBColor(0x0d, 0x0d, 0x0d), line_color=ACCENT_TEAL, line_width=Pt(1))
    add_text_box(slide, "$ npm install\n\n$ npm run dev\n\n-> Buka: http://localhost:5173",
                 left=0.55, top=2.35, width=3.9, height=1.4,
                 font_size=Pt(11), color=ACCENT_GREEN, font_name="Consolas")

    add_text_box(slide, "✅  Hot reload aktif\n✅  DevTools lengkap\n✅  Source maps tersedia",
                 left=0.45, top=3.9, width=4.1, height=0.9,
                 font_size=Pt(11), color=LIGHT_GRAY)

    # Opsi B
    add_rect(slide, 5.3, 1.35, 4.4, 3.8, fill_color=BG_CARD, line_color=ACCENT_BLUE, line_width=Pt(1.5))
    add_accent_bar(slide, 5.3, 1.35, 4.4, 0.06, ACCENT_BLUE)
    add_text_box(slide, "🌐  Opsi B — Standalone HTML", left=5.4, top=1.42, width=4.2, height=0.55,
                 font_size=Pt(14), bold=True, color=ACCENT_BLUE)
    add_text_box(slide, "Zero install — langsung buka di browser!",
                 left=5.4, top=1.9, width=4.1, height=0.35,
                 font_size=Pt(10), color=LIGHT_GRAY)

    add_text_box(slide, "Cukup klik dua kali:\nmqtt-chat-easy.html",
                 left=5.45, top=2.45, width=4.1, height=1.2,
                 font_size=Pt(18), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "✅  Tanpa Node.js\n✅  Tanpa npm install\n✅  Jalan di Chrome/Edge/Firefox",
                 left=5.45, top=3.9, width=4.1, height=0.9,
                 font_size=Pt(11), color=LIGHT_GRAY)

    # Footer
    add_rect(slide, 0.3, 5.35, 9.4, 1.15, fill_color=RGBColor(0x00, 0x3d, 0x33), line_color=ACCENT_GREEN, line_width=Pt(1))
    add_text_box(slide, "📡  Cara Connect Antar Laptop:", left=0.5, top=5.38, width=4, height=0.35,
                 font_size=Pt(11), bold=True, color=ACCENT_GREEN)
    add_text_box(slide, "Pastikan internet aktif  ->  Buka app di masing-masing laptop  ->  Klik Settings  ->  Isi username unik  ->  Pilih grup sama  ->  Mulai chat, foto, VN, video call!",
                 left=0.5, top=5.72, width=9.2, height=0.65,
                 font_size=Pt(10.5), color=LIGHT_GRAY)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 7: DEMO UI ─────────────────────────────────────────────────────────
def slide_ui_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "🖥️  Tampilan UI — Demo Aplikasi", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "Desain Glassmorphism Premium Dark Mode", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    # Sidebar
    add_rect(slide, 0.3, 1.35, 2.6, 5.1, fill_color=RGBColor(0x11, 0x1b, 0x21), line_color=ACCENT_TEAL, line_width=Pt(1))
    add_rect(slide, 0.3, 1.35, 2.6, 0.55, fill_color=ACCENT_TEAL)
    add_text_box(slide, "MQTT Chat", left=0.35, top=1.37, width=1.8, height=0.45,
                 font_size=Pt(11), bold=True, color=WHITE)
    add_text_box(slide, "⚙️", left=2.5, top=1.37, width=0.35, height=0.45,
                 font_size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)

    rooms = ["# General Lounge", "# Tech Talk", "# Random", "# Study Room", "# Gaming"]
    for i, room in enumerate(rooms):
        bg_col = RGBColor(0x00, 0x5c, 0x4b) if i == 0 else RGBColor(0x11, 0x1b, 0x21)
        add_rect(slide, 0.3, 1.92 + i * 0.7, 2.6, 0.65, fill_color=bg_col)
        add_text_box(slide, room, left=0.4, top=1.95 + i * 0.7, width=2.4, height=0.4,
                     font_size=Pt(10), color=WHITE if i == 0 else LIGHT_GRAY)
        if i == 0:
            add_text_box(slide, "🟢 5 online", left=0.4, top=2.22, width=2.3, height=0.25,
                         font_size=Pt(8), color=ACCENT_GREEN)

    # Chat area
    add_rect(slide, 2.95, 1.35, 6.75, 5.1, fill_color=RGBColor(0x0b, 0x14, 0x1a), line_color=ACCENT_TEAL, line_width=Pt(1))
    add_rect(slide, 2.95, 1.35, 6.75, 0.55, fill_color=RGBColor(0x1f, 0x2c, 0x34))
    add_text_box(slide, "# General Lounge", left=3.05, top=1.37, width=3, height=0.45,
                 font_size=Pt(12), bold=True, color=WHITE)
    add_text_box(slide, "📞  📹", left=8.9, top=1.37, width=0.75, height=0.45,
                 font_size=Pt(13), color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    bubbles = [
        (False, "Rian", "Halo semuanya! 👋", "10:30"),
        (True,  "Kamu", "Halo Rian! Gimana kabarnya?", "10:31"),
        (False, "Siti", "📷  [Foto Pantai.jpg]", "10:32"),
        (True,  "Kamu", "Wah keren banget fotonya! 🔥", "10:33"),
        (False, "Rian", "🎙️  Voice Note (0:12)", "10:34"),
    ]
    for i, (is_out, name, msg, time) in enumerate(bubbles):
        y_pos = 2.0 + i * 0.7
        if is_out:
            add_rect(slide, 6.2, y_pos, 3.3, 0.55, fill_color=RGBColor(0x00, 0x5c, 0x4b))
            add_text_box(slide, msg, left=6.3, top=y_pos+0.04, width=2.9, height=0.35,
                         font_size=Pt(9.5), color=WHITE)
            add_text_box(slide, time + " ✓✓", left=6.3, top=y_pos+0.35, width=3.0, height=0.2,
                         font_size=Pt(7.5), color=ACCENT_GREEN, align=PP_ALIGN.RIGHT)
        else:
            add_rect(slide, 3.05, y_pos, 3.3, 0.55, fill_color=RGBColor(0x1f, 0x2c, 0x34))
            add_text_box(slide, name + ":", left=3.15, top=y_pos+0.01, width=2, height=0.2,
                         font_size=Pt(7.5), bold=True, color=ACCENT_GREEN)
            add_text_box(slide, msg, left=3.15, top=y_pos+0.18, width=2.9, height=0.3,
                         font_size=Pt(9.5), color=WHITE)
            add_text_box(slide, time, left=5.0, top=y_pos+0.35, width=1.0, height=0.2,
                         font_size=Pt(7.5), color=LIGHT_GRAY)

    add_rect(slide, 2.95, 5.7, 6.75, 0.28, fill_color=RGBColor(0x1f, 0x2c, 0x34))
    add_text_box(slide, "✍️  Rian sedang mengetik...", left=3.05, top=5.72, width=5, height=0.23,
                 font_size=Pt(9), color=ACCENT_GREEN)

    add_rect(slide, 2.95, 6.0, 6.75, 0.45, fill_color=RGBColor(0x1f, 0x2c, 0x34), line_color=ACCENT_TEAL, line_width=Pt(0.5))
    add_text_box(slide, "📎  🖼️    Ketik pesan...                                    🎙️  ➤",
                 left=3.05, top=6.03, width=6.55, height=0.35,
                 font_size=Pt(9.5), color=LIGHT_GRAY)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 8: TEKNOLOGI & LIBRARY ─────────────────────────────────────────────
def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "📦  Teknologi & Library", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "Stack lengkap yang digunakan dalam proyek", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    techs = [
        ("⚛️",  "React 18",           "Framework UI berbasis komponen dengan hooks (useState, useEffect, useRef) untuk manajemen state real-time.",    ACCENT_GREEN, 0, 0),
        ("⚡",  "Vite 5",             "Build tool & dev server ultra-cepat dengan HMR dan bundling ES Module native.",                                  ACCENT_BLUE,  0, 1),
        ("🎨",  "Tailwind CSS 3",     "Utility-first CSS. Konfigurasi warna custom WhatsApp gelap (chat-bg, chat-accent, dll).",                        YELLOW,       0, 2),
        ("📡",  "MQTT.js 5",          "Client MQTT JavaScript. Koneksi HiveMQ via WSS. Publish/subscribe pesan, presence, dan sinyal WebRTC.",          ORANGE,       1, 0),
        ("🎥",  "WebRTC (Native)",    "API browser untuk P2P audio/video. RTCPeerConnection + STUN server Google.",                                     ACCENT_GREEN, 1, 1),
        ("🖼️", "Canvas API",         "Kompresi gambar sisi klien. Resize & re-render ke canvas sebelum kirim agar ukuran < 150 KB.",                   ACCENT_BLUE,  1, 2),
        ("🎙️", "MediaRecorder API",  "Rekam audio dari mikrofon secara native. Output Blob WebM dikonversi ke Base64.",                                YELLOW,       2, 0),
        ("🎇",  "Canvas Confetti",    "Efek kembang api diaktifkan dengan perintah /confetti di chat.",                                                 ORANGE,       2, 1),
        ("🔣",  "Lucide React",       "Koleksi 1000+ ikon SVG React. Dipakai untuk ikon telepon, kamera, mic, send, dll.",                             ACCENT_GREEN, 2, 2),
    ]

    card_w, card_h = 2.95, 1.48
    for (icon, name, desc, color, row, col) in techs:
        x = 0.25 + col * 3.2
        y = 1.38 + row * 1.6

        add_rect(slide, x, y, card_w, card_h, fill_color=BG_CARD, line_color=color, line_width=Pt(1))
        add_accent_bar(slide, x, y, card_w, 0.055, color)
        add_text_box(slide, f"{icon}  {name}", left=x+0.1, top=y+0.08, width=card_w-0.15, height=0.38,
                     font_size=Pt(12), bold=True, color=color)
        add_text_box(slide, desc, left=x+0.1, top=y+0.44, width=card_w-0.15, height=0.95,
                     font_size=Pt(9), color=LIGHT_GRAY)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 9: TANTANGAN & SOLUSI ──────────────────────────────────────────────
def slide_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    add_rect(slide, 0, 0.08, 10, 1.1, fill_color=BG_CARD)
    add_text_box(slide, "⚡  Tantangan & Solusi", left=0.3, top=0.15, width=9, height=0.7,
                 font_size=Pt(28), bold=True, color=WHITE)
    add_text_box(slide, "Masalah teknis yang dihadapi & cara mengatasinya", left=0.3, top=0.75, width=9, height=0.35,
                 font_size=Pt(13), color=ACCENT_GREEN)

    challenges = [
        ("Limit Payload Broker MQTT",
         "MQTT broker memiliki batas ukuran payload sekitar 256 KB. Gambar 10 MB tidak bisa dikirim langsung.",
         "Kompresi Canvas API sisi klien: gambar di-resize & kualitas dikurangi hingga <150 KB sebelum encode Base64.",
         ORANGE),
        ("Signaling WebRTC Tanpa Server Dedicated",
         "WebRTC butuh pertukaran SDP dan ICE candidates antar peer sebelum koneksi P2P bisa dibuka.",
         "MQTT dimanfaatkan sebagai signaling channel — topik unik per user untuk routing offer/answer/ICE candidates.",
         ACCENT_BLUE),
        ("Voice Note Format & Kompatibilitas Browser",
         "Format audio berbeda antar browser (Chrome=WebM/Opus, Firefox=OGG). Pengiriman biner langsung terlalu besar.",
         "MediaRecorder dengan MIME fallback + encode output Blob ke Base64 string sebelum publish ke MQTT topic.",
         YELLOW),
        ("Portabilitas: Butuh Node.js untuk React",
         "Tidak semua komputer memiliki Node.js terinstal — menyulitkan demo di laptop orang lain.",
         "Versi standalone mqtt-chat-easy.html yang memuat semua library via CDN — zero install, buka langsung di browser.",
         ACCENT_GREEN),
    ]

    for i, (title, problem, solution, color) in enumerate(challenges):
        y = 1.38 + i * 1.35
        add_rect(slide, 0.3, y, 9.4, 1.22, fill_color=BG_CARD, line_color=color, line_width=Pt(1))
        add_rect(slide, 0.3, y, 0.08, 1.22, fill_color=color)

        add_text_box(slide, "❌  " + title, left=0.5, top=y+0.04, width=9, height=0.32,
                     font_size=Pt(11), bold=True, color=color)
        add_text_box(slide, "Masalah: " + problem, left=0.5, top=y+0.34, width=9, height=0.35,
                     font_size=Pt(9.5), color=LIGHT_GRAY)
        add_text_box(slide, "✅ Solusi: " + solution, left=0.5, top=y+0.68, width=9, height=0.45,
                     font_size=Pt(9.5), color=WHITE)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── SLIDE 10: PENUTUP ────────────────────────────────────────────────────────
def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, 10, 0.08, ACCENT_GREEN)

    card = add_rect(slide, 0.5, 0.8, 9, 5.8, fill_color=BG_CARD, line_color=ACCENT_GREEN, line_width=Pt(1.5))
    add_accent_bar(slide, 0.5, 0.8, 9, 0.07, ACCENT_GREEN)

    add_text_box(slide, "🎉", left=3.8, top=1.1, width=2.4, height=1.2,
                 font_size=Pt(60), align=PP_ALIGN.CENTER)

    add_text_box(slide, "Terima Kasih!", left=0.5, top=2.2, width=9, height=0.85,
                 font_size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "MQTT WhatsApp Clone — Built From Scratch with ❤️", left=0.5, top=3.0, width=9, height=0.5,
                 font_size=Pt(16), color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    add_accent_bar(slide, 2.5, 3.5, 5, 0.04, ACCENT_TEAL)

    add_text_box(slide,
                 "✅  Real-time chat via MQTT (WSS) — latensi sub-detik\n"
                 "✅  Voice & Video Call P2P via WebRTC\n"
                 "✅  Voice Notes, Kompresi Foto 10MB, Typing & Presence Indicator\n"
                 "✅  Dua versi: React/Vite + Standalone HTML (zero install)\n"
                 "✅  Desain Glassmorphism Premium Dark Mode",
                 left=1.5, top=3.65, width=7, height=2.0,
                 font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_accent_bar(slide, 0, 6.8, 10, 0.08, ACCENT_GREEN)


# ─── MAIN ─────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7)

    print("🎨 Membuat slide 1: Cover...")
    slide_cover(prs)
    print("📋 Membuat slide 2: Overview Proyek...")
    slide_overview(prs)
    print("✨ Membuat slide 3: Fitur Utama...")
    slide_features(prs)
    print("🏗️  Membuat slide 4: Arsitektur Sistem...")
    slide_architecture(prs)
    print("🛠️  Membuat slide 5: Proses Development...")
    slide_development(prs)
    print("🚀 Membuat slide 6: Cara Menjalankan...")
    slide_howto(prs)
    print("🖥️  Membuat slide 7: Demo UI...")
    slide_ui_demo(prs)
    print("📦 Membuat slide 8: Teknologi & Library...")
    slide_tech(prs)
    print("⚡ Membuat slide 9: Tantangan & Solusi...")
    slide_challenges(prs)
    print("🎉 Membuat slide 10: Penutup...")
    slide_closing(prs)

    output = r"c:\Users\user\Documents\CHATTING\Presentasi_MQTT_WhatsApp_Clone.pptx"
    prs.save(output)
    print(f"\n✅ Presentasi berhasil disimpan!")
    print(f"📄 File: {output}")
    print(f"📊 Total slide: {len(prs.slides)}")


if __name__ == "__main__":
    main()

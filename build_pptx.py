from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG=RGBColor(0x0B,0x14,0x1A); SIDEBAR=RGBColor(0x11,0x1B,0x21)
ACCENT=RGBColor(0x00,0xA8,0x84); ACCENT_D=RGBColor(0x00,0x5C,0x4B)
WHITE=RGBColor(0xFF,0xFF,0xFF); LIGHT=RGBColor(0xE9,0xED,0xE9)
GRAY=RGBColor(0x86,0x96,0x8E)

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def bg(s,c=BG):
    b=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    b.line.fill.background(); b.fill.solid(); b.fill.fore_color.rgb=c; b.shadow.inherit=False
    return b

def tx(s,l,t,w,h,text,size=18,color=WHITE,bold=False,align=PP_ALIGN.LEFT,font="Segoe UI"):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=MSO_ANCHOR.TOP
    for i,ln in enumerate(str(text).split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name=font
    return tb

def bar(s,l,t,w=Inches(0.12),h=Inches(0.5)):
    b=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    b.line.fill.background(); b.fill.solid(); b.fill.fore_color.rgb=ACCENT; b.shadow.inherit=False
    return b

def title(s,t,sub=None):
    bar(s,Inches(0.6),Inches(0.55),h=Inches(0.6))
    tx(s,Inches(0.85),Inches(0.45),Inches(11),Inches(0.7),t,size=34,bold=True)
    if sub: tx(s,Inches(0.85),Inches(1.05),Inches(11),Inches(0.4),sub,size=16,color=ACCENT)

def card(s,l,t,w,h,c=SIDEBAR):
    o=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    o.line.fill.background(); o.fill.solid(); o.fill.fore_color.rgb=c; o.shadow.inherit=False
    try: o.adjustments[0]=0.06
    except: pass
    return o

def circ(s,l,t,size=Inches(0.45),text="",c=ACCENT,tc=WHITE,ts=16):
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,size,size)
    d.line.fill.background(); d.fill.solid(); d.fill.fore_color.rgb=c; d.shadow.inherit=False
    if text:
        tf=d.text_frame; tf.margin_top=Inches(0); tf.margin_bottom=Inches(0)
        tf.margin_left=Inches(0); tf.margin_right=Inches(0)
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=text; r.font.size=Pt(ts); r.font.bold=True; r.font.color.rgb=tc
    return d

I=Inches

# SLIDE 1: TITLE
s=prs.slides.add_slide(BLANK); bg(s)
t=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,I(0.25))
t.line.fill.background(); t.fill.solid(); t.fill.fore_color.rgb=ACCENT; t.shadow.inherit=False
bar(s,I(1),I(2.3),w=I(1.2),h=I(0.8))
tx(s,I(1),I(2.1),I(11),I(1),"Portfolio Project Antigravity",size=46,bold=True)
tx(s,I(1),I(3.3),I(11),I(0.6),"Game 3D FPS & WhatsApp Chat Clone Real-Time",size=22,color=ACCENT)
tx(s,I(1),I(4.0),I(11),I(0.6),"Dibuat dari awal sampai berjalan - Full Journey",size=16,color=LIGHT)
tx(s,I(1),I(6.4),I(11),I(0.4),"CrossFire 3D  |  MQTT WhatsApp Clone  |  React + Vite  |  Three.js  |  WebRTC",size=13,color=GRAY)

# SLIDE 2: AGENDA
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Agenda Presentasi","Apa yang akan dibahas")
items=[
 ("1","Proyek 1: CrossFire 3D","Game FPS 3D berbasis Three.js + Cannon-es"),
 ("2","Proyek 2: WhatsApp Clone","Chat real-time MQTT + WebRTC"),
 ("3","Stack Teknologi","Tools yang dipakai kedua proyek"),
 ("4","Langkah Pembuatan","Dari init sampai jalan"),
 ("5","Fitur Utama","Kemampuan setiap aplikasi"),
 ("6","Cara Menjalankan","Demo & penggunaan akhir"),
]
y=2.3
for n,t_,d in items:
    card(s,I(0.85),I(y),I(11.6),I(0.7))
    circ(s,I(1.05),I(y+0.13),text=n)
    tx(s,I(1.7),I(y+0.06),I(5),I(0.3),t_,size=15,bold=True)
    tx(s,I(1.7),I(y+0.35),I(9.5),I(0.3),d,size=12,color=GRAY)
    y+=0.78

# SLIDE 3: OVERVIEW CROSSFIRE
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Proyek 1: CrossFire 3D","3D Modern Warfare Browser Game")
card(s,I(0.85),I(1.7),I(11.6),I(5.0))
tx(s,I(1.1),I(1.9),I(11),I(0.4),"Game FPS (First Person Shooter) 3D yang berjalan langsung di browser.",size=17,bold=True,color=ACCENT)
tx(s,I(1.1),I(2.5),I(11),I(4),
   "Dibangun dengan Three.js untuk grafik 3D dan Cannon-es untuk fisika realistis.\n\n"
   "Fitur utama:\n"
   "- Multiplayer real-time via WebSocket (ws)\n"
   "- Server Express statis menyajikan folder public/\n"
   "- Engine fisika: gravitasi, collision, gerak pemain\n"
   "- Grafik WebGL 60 FPS langsung di browser\n"
   "- ID pemain otomatis di-generate server-side\n\n"
   "File kunci: server.js (multiplayer + static), package.json (deps).\n"
   "Dibuat menggunakan Antigravity dari nol hingga siap main.",
   size=13,color=LIGHT)

# SLIDE 4: OVERVIEW CHAT
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Proyek 2: WhatsApp Web Clone","Real-time chat dengan MQTT + WebRTC")
card(s,I(0.85),I(1.7),I(11.6),I(5.0))
tx(s,I(1.1),I(1.9),I(11),I(0.4),"Aplikasi chat bergaya WhatsApp Web bertema gelap premium.",size=17,bold=True,color=ACCENT)
tx(s,I(1.1),I(2.5),I(11),I(4),
   "Dibangun dengan React + Vite, protokol MQTT lewat WebSocket untuk pesan instan.\n\n"
   "Fitur utama:\n"
   "- Pesan instan sub-detik via MQTT (wss:// broker HiveMQ)\n"
   "- Voice & Video Call P2P pakai WebRTC (MQTT sebagai signaling)\n"
   "- Voice notes asli via MediaRecorder API (Base64 WebM)\n"
   "- Kirim gambar hingga 10MB (kompresi Canvas API)\n"
   "- Presence & typing indicator dinamis\n"
   "- Versi standalone: mqtt-chat-easy.html (tanpa instalasi)\n\n"
   "Dibuat menggunakan Antigravity dari nol hingga bisa digunakan.",
   size=13,color=LIGHT)

# SLIDE 5: STACK TEKNOLOGI
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Stack Teknologi","Tools kedua proyek")
techs=[
 ("Three.js","Grafik 3D WebGL (CrossFire)",ACCENT),
 ("Cannon-es","Engine fisika 3D (CrossFire)",ACCENT_D),
 ("Express","Server statis + routing",ACCENT),
 ("WebSocket","Multiplayer real-time game",ACCENT_D),
 ("React 18","UI komponen (Chat)",ACCENT),
 ("Vite","Dev server super cepat",ACCENT_D),
 ("MQTT.js","Messaging chat real-time",ACCENT),
 ("WebRTC","Voice/video call P2P",ACCENT_D),
]
cols=4; cw,ch=I(2.8),I(2.2); gx,gy=I(0.35),I(0.35)
sx,sy=I(0.85),I(1.8)
for i,(n,d,c) in enumerate(techs):
    r=i//cols; col=i%cols
    x=sx+col*(cw+gx); y=sy+r*(ch+gy)
    card(s,x,y,cw,ch)
    st=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,cw,I(0.15))
    st.line.fill.background(); st.fill.solid(); st.fill.fore_color.rgb=c; st.shadow.inherit=False
    try: st.adjustments[0]=0.5
    except: pass
    tx(s,x,y+I(0.5),cw,I(0.5),n,size=18,bold=True,align=PP_ALIGN.CENTER)
    tx(s,x+I(0.15),y+I(1.05),cw-I(0.3),I(1),d,size=12,color=GRAY,align=PP_ALIGN.CENTER)

# SLIDE 6: LANGKAH CROSSFIRE
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Langkah Buat CrossFire 3D","Dari init sampai jalan")
steps=[
 ("1","Init Proyek","npm init, install three, cannon-es, express, ws"),
 ("2","Buat Server","server.js: Express static + WebSocket port 3000"),
 ("3","Multiplayer Logic","Generate clientId, broadcast posisi via ws"),
 ("4","Buat public/","index.html + client.js render 3D + input"),
 ("5","Integrasi Fisika","Cannon-es: gravitasi, collision pemain"),
 ("6","Jalankan","npm start -> buka http://localhost:3000"),
]
y=1.7
for i,(n,t_,d) in enumerate(steps):
    r=i//2; col=i%2
    x=I(0.85)+col*I(6.2); yy=I(y+r*1.6)
    card(s,x,yy,I(5.9),I(1.4))
    circ(s,x+I(0.2),yy+I(0.45),text=n)
    tx(s,x+I(0.8),yy+I(0.15),I(5),I(0.4),t_,size=16,bold=True)
    tx(s,x+I(0.8),yy+I(0.6),I(5),I(0.7),d,size=12,color=GRAY)

# SLIDE 7: LANGKAH CHAT
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Langkah Buat WhatsApp Clone","Dari init sampai bisa chat")
steps=[
 ("1","Init Vite+React","npm create vite CHATTING, npm install"),
 ("2","Tailwind CSS","Palet WhatsApp dark, animasi custom"),
 ("3","Install Deps","mqtt, lucide-react, canvas-confetti, aedes"),
 ("4","Logika App.jsx","Koneksi MQTT, presence, typing, VN"),
 ("5","WebRTC Call","createOffer -> signaling MQTT -> ICE"),
 ("6","Standalone HTML","mqtt-chat-easy.html + git commit awal"),
]
y=1.7
for i,(n,t_,d) in enumerate(steps):
    r=i//2; col=i%2
    x=I(0.85)+col*I(6.2); yy=I(y+r*1.6)
    card(s,x,yy,I(5.9),I(1.4))
    circ(s,x+I(0.2),yy+I(0.45),text=n)
    tx(s,x+I(0.8),yy+I(0.15),I(5),I(0.4),t_,size=16,bold=True)
    tx(s,x+I(0.8),yy+I(0.6),I(5),I(0.7),d,size=12,color=GRAY)

# SLIDE 8: FITUR CROSSFIRE
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Fitur CrossFire 3D","Kemampuan game FPS browser")
feats=[
 ("Grafik 3D WebGL","Render Three.js 60 FPS langsung di browser"),
 ("Multiplayer Real-Time","WebSocket sinkronisasi posisi pemain"),
 ("Fisika Realistis","Cannon-es: gravitasi, collision, gerak"),
 ("First Person Control","Mouse look + WASD movement modern"),
 ("Server Statis Express","Sajikan public/ dan node_modules"),
 ("Auto Player ID","clientId di-generate server, broadcast init"),
]
y=1.7
for i,(t_,d) in enumerate(feats):
    r=i//2; col=i%2
    x=I(0.85)+col*I(6.2); yy=I(y+r*1.55)
    card(s,x,yy,I(5.9),I(1.35))
    circ(s,x+I(0.2),yy+I(0.45),text=str(i+1))
    tx(s,x+I(0.8),yy+I(0.15),I(5),I(0.4),t_,size=16,bold=True)
    tx(s,x+I(0.8),yy+I(0.6),I(5),I(0.7),d,size=12,color=GRAY)

# SLIDE 9: FITUR CHAT
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Fitur WhatsApp Clone","Kemampuan chat real-time")
feats=[
 ("Pesan Instan","MQTT over WebSocket wss:// sub-detik"),
 ("Voice & Video Call","WebRTC P2P, MQTT signaling channel"),
 ("Voice Notes Asli","MediaRecorder API, Base64 WebM <250KB"),
 ("Kompresi Gambar 10MB","Canvas API, output <150KB sisi klien"),
 ("Presence & Typing","User online & status mengetik dinamis"),
 ("Glassmorphism Premium","Tema gelap, animasi, responsif"),
]
y=1.7
for i,(t_,d) in enumerate(feats):
    r=i//2; col=i%2
    x=I(0.85)+col*I(6.2); yy=I(y+r*1.55)
    card(s,x,yy,I(5.9),I(1.35))
    circ(s,x+I(0.2),yy+I(0.45),text=str(i+1))
    tx(s,x+I(0.8),yy+I(0.15),I(5),I(0.4),t_,size=16,bold=True)
    tx(s,x+I(0.8),yy+I(0.6),I(5),I(0.7),d,size=12,color=GRAY)

# SLIDE 10: CARA JALANKAN
s=prs.slides.add_slide(BLANK); bg(s); title(s,"Cara Menjalankan","Dua proyek siap pakai")
card(s,I(0.85),I(1.7),I(5.8),I(5.0))
tx(s,I(1.05),I(1.85),I(5.4),I(0.4),"CrossFire 3D",size=18,bold=True,color=ACCENT)
tx(s,I(1.05),I(2.4),I(5.4),I(4),
   "Jalankan game FPS di lokal:\n\n"
   "$ cd Documents/CrossFire\n"
   "$ npm install\n"
   "$ npm start\n\n"
   "Buka browser:\n"
   "http://localhost:3000\n\n"
   "Game langsung siap dimainkan!",
   size=13,color=LIGHT)
card(s,I(6.85),I(1.7),I(5.6),I(5.0))
tx(s,I(7.05),I(1.85),I(5.2),I(0.4),"WhatsApp Clone",size=18,bold=True,color=ACCENT)
tx(s,I(7.05),I(2.4),I(5.2),I(4),
   "Dua opsi menjalankan chat:\n\n"
   "Opsi A - Developer mode:\n"
   "$ cd Documents/CHATTING\n"
   "$ npm install\n"
   "$ npm run dev -> localhost:5173\n\n"
   "Opsi B - Gampang (tanpa Node.js):\n"
   "Klik ganda mqtt-chat-easy.html\n"
   "Langsung buka di browser!",
   size=13,color=LIGHT)

# SLIDE 11: PENUTUP
s=prs.slides.add_slide(BLANK); bg(s)
t=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,I(0.25))
t.line.fill.background(); t.fill.solid(); t.fill.fore_color.rgb=ACCENT; t.shadow.inherit=False
bar(s,I(1),I(2.7),w=I(1.2),h=I(0.8))
tx(s,I(1),I(2.5),I(11),I(1),"Terima Kasih",size=54,bold=True)
tx(s,I(1),I(3.7),I(11),I(0.6),"Dua proyek selesai dibuat dengan Antigravity",size=20,color=ACCENT)
tx(s,I(1),I(4.4),I(11),I(0.6),"CrossFire 3D FPS Game  &  MQTT WhatsApp Web Clone",size=16,color=LIGHT)
tx(s,I(1),I(6.6),I(11),I(0.4),"React + Vite  |  Three.js + Cannon-es  |  MQTT + WebRTC  |  Tailwind CSS",size=13,color=GRAY)

out=r"C:\Users\user\Documents\CHATTING\Portfolio_Antigravity.pptx"
prs.save(out)
print("SAVED:", out)
print("Total slides:", len(prs.slides._sldIdLst))

